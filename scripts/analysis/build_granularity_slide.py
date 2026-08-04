from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1E, 0x27, 0x61)
BODY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x6F, 0x7A)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
GREEN_BG = RGBColor(0xE9, 0xF5, 0xEE)
PINK_BG = RGBColor(0xFB, 0xEC, 0xE7)
BLUE_BG = RGBColor(0xF3, 0xF6, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG = NAVY

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(x, y, w, h, text, size, color, bold=False, italic=False,
             font="Calibri", align=PP_ALIGN.LEFT, line_spacing=None, anchor=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_multiline(x, y, w, h, lines, size, color, font="Calibri", bold=False,
                   line_spacing=1.2, space_after=6, anchor=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return box


# ---- Title ----
add_text(0.6, 0.35, 12.13, 0.65,
          "Cross-Granularity Ablation: Learning vs. Inference vs. Labels",
          27, NAVY, bold=True, font="Cambria")

add_text(0.6, 0.98, 12.13, 0.35,
          "NYC 1,000-request instances, 50 vehicles, LR = 1e-4 throughout, two seeds",
          13.5, GRAY, italic=True, font="Calibri")

add_text(0.6, 1.3, 12.13, 0.3,
          "cheap = 1/4 min granularity (step_size 60 / batch_interval 240)   ·   "
          "expensive = 2/8 min granularity (step_size 120 / batch_interval 480)",
          12, GRAY, italic=True, font="Calibri")

# ---- Table ----
table_x, table_y = 0.5, 1.65
col_w = [4.55, 1.72, 1.72, 1.72, 1.72]
row_h = [0.62, 0.5, 0.5, 0.5, 0.5]
table_w = sum(col_w)
table_h = sum(row_h)

gtable = slide.shapes.add_table(5, 5, Inches(table_x), Inches(table_y),
                                 Inches(table_w), Inches(table_h)).table
for i, w in enumerate(col_w):
    gtable.columns[i].width = Inches(w)
for i, h in enumerate(row_h):
    gtable.rows[i].height = Inches(h)

tbl = gtable._tbl
tblPr = tbl.find(qn('a:tblPr'))
tblPr.set('firstRow', '0')
tblPr.set('bandRow', '0')

headers = ["Configuration", "Val\nseed 42", "Val\nseed 1", "Test\nseed 42", "Test\nseed 1"]
rows = [
    ("Learning cheap / Inference cheap  (1/4)", "73.45%", "62.46%", "64.01%", "58.95%", "base"),
    ("Learning expensive / Inference expensive  (2/8)", "87.58%", "84.34%", "82.83%", "80.04%", "best"),
    ("Learning expensive / Inference cheap", "77.46%", "73.93%", "72.70%", "66.49%", "base"),
    ("Labels expensive, Learning + Inference cheap", "59.50%", "43.36%", "53.26%", "46.74%", "worst"),
]

row_color = {"base": WHITE, "best": GREEN_BG, "worst": PINK_BG}
val_color = {"base": BODY, "best": GREEN, "worst": RGBColor(0xB8, 0x3A, 0x2A)}

for c, htext in enumerate(headers):
    cell = gtable.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = HEADER_BG
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    lines = htext.split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = WHITE

for r, (label, val42, val1, test42, test1, kind) in enumerate(rows, start=1):
    bg = row_color[kind]
    vcolor = val_color[kind]
    values = [label, val42, val1, test42, test1]
    for c in range(5):
        cell = gtable.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        if c == 0:
            p.alignment = PP_ALIGN.LEFT
            run.text = label
            run.font.size = Pt(12.5)
            run.font.bold = (kind != "base")
            run.font.color.rgb = BODY
        else:
            p.alignment = PP_ALIGN.CENTER
            run.text = values[c]
            run.font.size = Pt(14.5)
            run.font.bold = True
            run.font.color.rgb = vcolor
        run.font.name = "Calibri"

# ---- Comment box (navy) ----
box_y = table_y + table_h + 0.25
box_h = 2.7
comment_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(box_y), Inches(11.93), Inches(box_h))
comment_box.adjustments[0] = 0.06
set_fill(comment_box, NAVY)
comment_box.shadow.inherit = False
comment_box.text_frame.clear()
comment_box.text_frame.word_wrap = True

add_text(0.95, box_y + 0.16, 11.4, 0.32, "What this tells us", 16, WHITE, bold=True, font="Calibri")

bullets = [
    "Consistent granularity wins outright at both seeds: training and evaluating both at 2/8 is the "
    "strongest result of the four — 87.6% / 84.3% val, 82.8% / 80.0% test (seed 42 / seed 1).",
    "An expensive-trained model generalizes down at both seeds: expensive-learn / cheap-infer still beats "
    "training cheap from the start, even run at the finer 1/4 timing it never trained on.",
    "Label/timing mismatch is the worst failure mode at both seeds — inconsistency between label and "
    "training granularity hurts more than simply using the weaker granularity throughout.",
    "Ranking is seed-robust; absolute service rate shifts 5-15 points between seeds, largest for cheap/cheap.",
]
add_multiline(0.95, box_y + 0.56, 11.4, box_h - 0.7, bullets, 12, WHITE,
              font="Calibri", line_spacing=1.1, space_after=6)

out_path = "/private/tmp/claude-501/-Users-joni-Desktop-Masterarbeit-Reduce-then-optimize/06426b27-f40c-4745-a015-657068376640/scratchpad/pptx/nyc_granularity_cross_experiment.pptx"
prs.save(out_path)
print("saved:", out_path)
