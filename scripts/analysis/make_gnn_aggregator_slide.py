"""Single slide explaining the three GNN aggregator variants (GCN/Mean/Pool)
used in the request-graph candidate scorer, built with python-pptx (no
Node.js available on this machine for pptxgenjs)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

TITLE_COLOR = RGBColor(0x21, 0x21, 0x21)
SUBTLE_GRAY = RGBColor(0x6B, 0x72, 0x80)
BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xDC)
NODE_SELF = RGBColor(0x36, 0x45, 0x4F)
NODE_NEIGHBOR = RGBColor(0xAD, 0xB5, 0xBD)
LINK_GRAY = RGBColor(0xB0, 0xB6, 0xBC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

VARIANTS = [
    {
        "code": "GCN",
        "title": "GCN — blend, then transform",
        "accent": RGBColor(0x1C, 0x72, 0x93),
        "bg": RGBColor(0xEA, 0xF3, 0xF5),
        "intuition": "“Average myself in with my neighbours, then react to the blend.”",
        "caption": "GCNMeanLayer (Eq. 2) — self + neighbours meaned into ONE vector, one shared linear layer.",
    },
    {
        "code": "MEAN",
        "title": "Mean — keep separate, then transform",
        "accent": RGBColor(0xC9, 0x7B, 0x1F),
        "bg": RGBColor(0xFB, 0xF1, 0xE4),
        "intuition": "“Keep my signal and my neighbours’ average distinct — let the network learn how to weigh them.”",
        "caption": "GraphSAGEMeanLayer (Alg. 1) — self and neighbour-mean concatenated, one shared linear layer.",
    },
    {
        "code": "POOL",
        "title": "Pool — pick the strongest signal",
        "accent": RGBColor(0xB2, 0x3A, 0x48),
        "bg": RGBColor(0xF7, 0xE9, 0xEB),
        "intuition": "“Each neighbour votes per feature — the strongest vote wins, instead of averaging.”",
        "caption": "GraphSAGEPoolLayer (Eq. 3) — neighbours transformed individually, element-wise MAX, then concat with self.",
    },
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def add_text(x, y, w, h, text, size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_oval(x, y, d, fill, label=None, label_color=WHITE, font_size=8):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    if label:
        tf = shp.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = label_color
        run.font.name = "Calibri"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def add_box(x, y, w, h, fill, label, label_color=WHITE, font_size=9, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = 0.18
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = label_color
    run.font.name = "Calibri"
    return shp


def add_diamond(x, y, w, h, fill, label, font_size=9):
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    return shp


def add_link(x1, y1, x2, y2, color=LINK_GRAY, width_pt=1.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    conn.shadow.inherit = False
    return conn


# ---- Title ----
add_text(0.5, 0.35, 12.3, 0.7, "Three GNN Aggregator Variants", 32, TITLE_COLOR, bold=True)
add_text(
    0.5, 1.05, 12.3, 0.4,
    "Combining a request node's own embedding with its neighbours in the conflict graph — three ways to do it.",
    14, SUBTLE_GRAY, italic=True,
)

# ---- Columns ----
col_w = 3.8
gap = 0.38
xs = [0.5, 0.5 + col_w + gap, 0.5 + 2 * (col_w + gap)]
card_y = 1.65
card_h = 5.35

for i, v in enumerate(VARIANTS):
    x0 = xs[i]
    accent = v["accent"]

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x0), Inches(card_y), Inches(col_w), Inches(card_h))
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = v["bg"]
    card.line.color.rgb = BORDER_GRAY
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    # Badge + title
    add_oval(x0 + 0.25, card_y + 0.25, 0.55, accent, label=v["code"], font_size=10.5)
    add_text(x0 + 0.95, card_y + 0.28, col_w - 1.2, 0.5, v["title"], 14.5, accent, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    # --- Input mini-graph (same structure for all three) ---
    g_top = card_y + 1.0
    n1_x, n2_x = x0 + 0.55, x0 + col_w - 0.55 - 0.4
    add_link(n1_x + 0.2, g_top + 0.4, x0 + col_w / 2, g_top + 0.75, LINK_GRAY, 1.0)
    add_link(n2_x + 0.2, g_top + 0.4, x0 + col_w / 2, g_top + 0.75, LINK_GRAY, 1.0)
    add_oval(n1_x, g_top, 0.4, NODE_NEIGHBOR, label="n", font_size=8)
    add_oval(n2_x, g_top, 0.4, NODE_NEIGHBOR, label="n", font_size=8)
    self_x = x0 + col_w / 2 - 0.25
    self_y = g_top + 0.65
    add_oval(self_x, self_y, 0.5, NODE_SELF, label="self", font_size=8)

    # --- Operation-specific diagram ---
    op_top = self_y + 0.65
    if v["code"] == "GCN":
        avg_x, avg_y, avg_d = x0 + col_w / 2 - 0.3, op_top, 0.6
        add_link(self_x + 0.25, self_y + 0.5, avg_x + 0.3, avg_y, accent, 1.2)
        add_link(n1_x + 0.2, g_top + 0.4, avg_x + 0.15, avg_y + 0.1, LINK_GRAY, 0.9)
        add_link(n2_x + 0.2, g_top + 0.4, avg_x + 0.45, avg_y + 0.1, LINK_GRAY, 0.9)
        add_oval(avg_x, avg_y, avg_d, accent, label="avg", font_size=9)
        lin_x, lin_y, lin_w, lin_h = x0 + col_w / 2 - 0.55, avg_y + 0.75, 1.1, 0.42
        add_link(avg_x + 0.3, avg_y + avg_d, lin_x + lin_w / 2, lin_y, accent, 1.2)
        add_box(lin_x, lin_y, lin_w, lin_h, accent, "Linear", font_size=9.5)
        out_y = lin_y + lin_h + 0.12
        add_text(x0, out_y, col_w, 0.3, "→ self′", 10.5, accent, bold=True, align=PP_ALIGN.CENTER)
    elif v["code"] == "MEAN":
        box_w, box_h = 1.35, 0.42
        self_box_x = x0 + 0.35
        avg_box_x = x0 + col_w - 0.35 - box_w
        box_y = op_top + 0.05
        add_link(self_x + 0.25, self_y + 0.5, self_box_x + box_w / 2, box_y, accent, 1.0)
        add_link(n1_x + 0.2, g_top + 0.4, avg_box_x + box_w / 2 - 0.2, box_y, LINK_GRAY, 0.9)
        add_link(n2_x + 0.2, g_top + 0.4, avg_box_x + box_w / 2 + 0.2, box_y, LINK_GRAY, 0.9)
        add_box(self_box_x, box_y, box_w, box_h, NODE_SELF, "self", font_size=9)
        add_box(avg_box_x, box_y, box_w, box_h, NODE_NEIGHBOR, "avg(nbrs)", label_color=RGBColor(0x21, 0x21, 0x21), font_size=8.5)
        add_text(x0, box_y + box_h * 0.5 - 0.12, col_w, 0.25, "∥", 14, accent, bold=True, align=PP_ALIGN.CENTER)
        lin_x, lin_y, lin_w, lin_h = x0 + col_w / 2 - 0.55, box_y + box_h + 0.28, 1.1, 0.42
        add_link(self_box_x + box_w / 2, box_y + box_h, lin_x + lin_w / 2 - 0.15, lin_y, accent, 1.0)
        add_link(avg_box_x + box_w / 2, box_y + box_h, lin_x + lin_w / 2 + 0.15, lin_y, accent, 1.0)
        add_box(lin_x, lin_y, lin_w, lin_h, accent, "Linear", font_size=9.5)
        out_y = lin_y + lin_h + 0.12
        add_text(x0, out_y, col_w, 0.3, "→ self′", 10.5, accent, bold=True, align=PP_ALIGN.CENTER)
    else:  # POOL
        w_w, w_h = 0.85, 0.36
        w1_x = n1_x - 0.22
        w2_x = n2_x - 0.22
        w_y = op_top - 0.05
        add_link(n1_x + 0.2, g_top + 0.4, w1_x + w_w / 2, w_y, LINK_GRAY, 0.9)
        add_link(n2_x + 0.2, g_top + 0.4, w2_x + w_w / 2, w_y, LINK_GRAY, 0.9)
        add_box(w1_x, w_y, w_w, w_h, accent, "W", font_size=9)
        add_box(w2_x, w_y, w_w, w_h, accent, "W", font_size=9)
        dia_x, dia_y, dia_w, dia_h = x0 + col_w / 2 - 0.35, w_y + w_h + 0.12, 0.7, 0.4
        add_link(w1_x + w_w / 2, w_y + w_h, dia_x + dia_w / 2 - 0.1, dia_y, accent, 1.0)
        add_link(w2_x + w_w / 2, w_y + w_h, dia_x + dia_w / 2 + 0.1, dia_y, accent, 1.0)
        add_diamond(dia_x, dia_y, dia_w, dia_h, accent, "max", font_size=9)
        lin_x, lin_y, lin_w, lin_h = x0 + col_w / 2 - 0.6, dia_y + dia_h + 0.12, 1.2, 0.36
        add_link(self_x + 0.25, self_y + 0.5, lin_x + lin_w / 2 - 0.25, lin_y, accent, 1.0)
        add_link(dia_x + dia_w / 2, dia_y + dia_h, lin_x + lin_w / 2 + 0.25, lin_y, accent, 1.0)
        add_box(lin_x, lin_y, lin_w, lin_h, accent, "Linear", font_size=9)
        out_y = lin_y + lin_h + 0.08
        add_text(x0, out_y, col_w, 0.3, "→ self′", 10.5, accent, bold=True, align=PP_ALIGN.CENTER)

    # Intuition + caption
    intuition_y = card_y + card_h - 1.55
    add_text(x0 + 0.28, intuition_y, col_w - 0.56, 0.85, v["intuition"], 12.5, TITLE_COLOR, italic=True)
    caption_y = card_y + card_h - 0.62
    add_text(x0 + 0.28, caption_y, col_w - 0.56, 0.5, v["caption"], 9.5, SUBTLE_GRAY)

out_path = "gnn_aggregator_variants_slide.pptx"
prs.save(out_path)
print("saved:", out_path)
